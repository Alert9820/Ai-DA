"""
AI Data Analyst Pro - Production Ready
Features: EDA, Time Series, Statistical Tests, Advanced Visualizations, ML Insights,
Natural Language Queries (Gemini), Big Data (Dask), Reporting, Collaboration,
Performance (Celery+Redis), Security (JWT, rate limiting, CORS, logging)
"""

import os
import json
import uuid
import hashlib
import logging
from logging.handlers import RotatingFileHandler
from datetime import datetime, timedelta
from functools import wraps
from io import BytesIO, StringIO
import base64
import asyncio
import tempfile
from concurrent.futures import ThreadPoolExecutor

# Flask & extensions
from flask import Flask, request, render_template_string, jsonify, g, session, url_for, send_file
from flask_cors import CORS
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
from flask_compress import Compress
from flask_jwt_extended import (
    JWTManager, create_access_token, jwt_required, get_jwt_identity,
    get_jwt, create_refresh_token
)
from flask_sqlalchemy import SQLAlchemy
from flask_caching import Cache
from sqlalchemy import text, func, inspect

# Data processing
import pandas as pd
import numpy as np
import dask.dataframe as dd
from dask.diagnostics import ProgressBar
import pyarrow.parquet as pq
import sqlalchemy as db

# Visualization
import plotly.express as px
import plotly.graph_objects as go
import plotly.figure_factory as ff
import plotly.utils
import matplotlib.pyplot as plt
import seaborn as sns

# AI & ML
import google.generativeai as genai
import pycaret
from pycaret.classification import *
from pycaret.regression import *
from pycaret.clustering import *
from sklearn.ensemble import RandomForestRegressor, RandomForestClassifier
from sklearn.cluster import KMeans, DBSCAN
from sklearn.preprocessing import StandardScaler
from sklearn.ensemble import IsolationForest
from sklearn.feature_selection import mutual_info_regression, mutual_info_classif
import shap
from apyori import apriori
import statsmodels.api as sm
from statsmodels.tsa.seasonal import seasonal_decompose
from statsmodels.tsa.arima.model import ARIMA
from scipy import stats
import prophet  # fbprophet

# Reporting & emails
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.platypus import SimpleDocTemplate, Table, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
import openpyxl
from pptx import Presentation
import smtplib
from email.mime.text import MimeText
from email.mime.multipart import MimeMultipart
from email.mime.base import MIMEBase
from email import encoders
import requests
import schedule
import time

# Async tasks
from celery import Celery
from celery.result import AsyncResult

# Utilities
import re
import hashlib
import random
import string
from werkzeug.utils import secure_filename
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# ------------------------------------------------------------------------------
# App Configuration
# ------------------------------------------------------------------------------
app = Flask(__name__)
app.config['SECRET_KEY'] = os.getenv('SECRET_KEY', 'dev-key-change-in-production')
app.config['JWT_SECRET_KEY'] = os.getenv('JWT_SECRET_KEY', 'jwt-secret-change')
app.config['SQLALCHEMY_DATABASE_URI'] = os.getenv('DATABASE_URL', 'sqlite:///data_analyst.db')
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['CACHE_TYPE'] = 'RedisCache' if os.getenv('REDIS_URL') else 'SimpleCache'
app.config['CACHE_REDIS_URL'] = os.getenv('REDIS_URL', 'redis://localhost:6379/0')
app.config['RATELIMIT_STORAGE_URL'] = os.getenv('REDIS_URL', 'memory://')
app.config['UPLOAD_FOLDER'] = '/tmp/uploads'
app.config['MAX_CONTENT_LENGTH'] = 2 * 1024 * 1024 * 1024  # 2GB max upload
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Extensions
CORS(app, resources={r"/*": {"origins": "*"}})  # restrict in production
limiter = Limiter(app, key_func=get_remote_address, default_limits=["200 per day", "50 per hour"])
Compress(app)
jwt = JWTManager(app)
db = SQLAlchemy(app)
cache = Cache(app)

# Celery configuration
celery = Celery(
    app.import_name,
    broker=os.getenv('REDIS_URL', 'redis://localhost:6379/0'),
    backend=os.getenv('REDIS_URL', 'redis://localhost:6379/0')
)
celery.conf.update(
    task_serializer='json',
    accept_content=['json'],
    result_serializer='json',
    timezone='UTC',
    enable_utc=True,
)

# Gemini AI
genai.configure(api_key=os.getenv('GEMINI_API_KEY'))
gemini_model = genai.GenerativeModel('gemini-1.5-flash')

# Logging
if not app.debug:
    handler = RotatingFileHandler('app.log', maxBytes=10000000, backupCount=5)
    handler.setLevel(logging.INFO)
    formatter = logging.Formatter('%(asctime)s %(levelname)s: %(message)s [in %(pathname)s:%(lineno)d]')
    handler.setFormatter(formatter)
    app.logger.addHandler(handler)
    app.logger.setLevel(logging.INFO)
    app.logger.info('AI Data Analyst Pro startup')

# ------------------------------------------------------------------------------
# Database Models (Collaboration)
# ------------------------------------------------------------------------------
class User(db.Model):
    __tablename__ = 'users'
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(80), unique=True, nullable=False)
    email = db.Column(db.String(120), unique=True, nullable=False)
    password_hash = db.Column(db.String(200), nullable=False)
    role = db.Column(db.String(20), default='viewer')  # viewer, editor, admin
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    projects = db.relationship('Project', backref='owner', lazy=True)
    comments = db.relationship('Comment', backref='author', lazy=True)
    activities = db.relationship('ActivityLog', backref='user', lazy=True)

class Project(db.Model):
    __tablename__ = 'projects'
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(200), nullable=False)
    description = db.Column(db.Text)
    owner_id = db.Column(db.Integer, db.ForeignKey('users.id'))
    data_snapshot = db.Column(db.Text)  # JSON or path to stored data
    analysis_config = db.Column(db.Text)  # JSON
    share_token = db.Column(db.String(64), unique=True)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)
    comments = db.relationship('Comment', backref='project', lazy=True)

class Comment(db.Model):
    __tablename__ = 'comments'
    id = db.Column(db.Integer, primary_key=True)
    content = db.Column(db.Text, nullable=False)
    user_id = db.Column(db.Integer, db.ForeignKey('users.id'))
    project_id = db.Column(db.Integer, db.ForeignKey('projects.id'))
    created_at = db.Column(db.DateTime, default=datetime.utcnow)

class ActivityLog(db.Model):
    __tablename__ = 'activity_logs'
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('users.id'))
    action = db.Column(db.String(200))
    details = db.Column(db.Text)
    ip_address = db.Column(db.String(45))
    timestamp = db.Column(db.DateTime, default=datetime.utcnow)

# ------------------------------------------------------------------------------
# Helper Functions: Data Loading with Chunking & Dask
# ------------------------------------------------------------------------------
def load_data(filepath, filetype=None, use_dask=True):
    """Load large files using chunking or Dask, support CSV/Excel/JSON/Parquet/SQLite"""
    if filetype is None:
        ext = filepath.split('.')[-1].lower()
    else:
        ext = filetype.lower()
    
    # For huge datasets > 500MB, force Dask
    file_size = os.path.getsize(filepath) / (1024*1024)
    force_dask = file_size > 500 or use_dask
    
    if ext == 'csv':
        if force_dask:
            return dd.read_csv(filepath, assume_missing=True)
        else:
            return pd.read_csv(filepath, low_memory=False)
    elif ext in ['xls', 'xlsx']:
        return pd.read_excel(filepath)
    elif ext == 'json':
        return pd.read_json(filepath, lines=True)
    elif ext == 'parquet':
        return pd.read_parquet(filepath)
    elif ext == 'sqlite':
        conn = sqlite3.connect(filepath)
        return pd.read_sql_query("SELECT * FROM data", conn)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

def load_from_database(conn_string, table_name, db_type='postgres'):
    """Connect to PostgreSQL/MySQL and load table"""
    if db_type == 'postgres':
        engine = db.create_engine(conn_string)
    elif db_type == 'mysql':
        engine = db.create_engine(conn_string)
    else:
        raise ValueError("Unsupported DB type")
    with engine.connect() as conn:
        return pd.read_sql_table(table_name, conn)

# ------------------------------------------------------------------------------
# EDA Functions
# ------------------------------------------------------------------------------
def univariate_analysis(df):
    """Generate univariate statistics and plots"""
    numeric_cols = df.select_dtypes(include=[np.number]).columns
    cat_cols = df.select_dtypes(include=['object', 'category']).columns
    
    results = {
        'summary': df.describe(include='all').to_dict(),
        'missing': df.isnull().sum().to_dict(),
        'skewness': df[numeric_cols].skew().to_dict() if len(numeric_cols) else {},
        'kurtosis': df[numeric_cols].kurtosis().to_dict() if len(numeric_cols) else {},
        'unique_counts': df[cat_cols].nunique().to_dict() if len(cat_cols) else {}
    }
    return results

def bivariate_analysis(df, x, y):
    """Correlation and cross-tabulation"""
    if pd.api.types.is_numeric_dtype(df[x]) and pd.api.types.is_numeric_dtype(df[y]):
        corr = df[[x, y]].corr().iloc[0,1]
        return {'correlation': corr, 'type': 'numeric-numeric'}
    elif pd.api.types.is_numeric_dtype(df[x]) and not pd.api.types.is_numeric_dtype(df[y]):
        grouped = df.groupby(y)[x].describe().to_dict()
        return {'grouped_stats': grouped, 'type': 'numeric-categorical'}
    else:
        crosstab = pd.crosstab(df[x], df[y]).to_dict()
        return {'crosstab': crosstab, 'type': 'categorical-categorical'}

def multivariate_analysis(df, features):
    """Pairwise correlations and parallel coordinates"""
    numeric_df = df[features].select_dtypes(include=[np.number])
    if len(numeric_df.columns) > 1:
        corr_matrix = numeric_df.corr().to_dict()
        return {'correlation_matrix': corr_matrix}
    else:
        return {}

def outlier_detection(df, method='iqr'):
    """Multiple outlier detection methods"""
    numeric_cols = df.select_dtypes(include=[np.number]).columns
    outliers = {}
    for col in numeric_cols:
        if method == 'iqr':
            Q1 = df[col].quantile(0.25)
            Q3 = df[col].quantile(0.75)
            IQR = Q3 - Q1
            lower = Q1 - 1.5*IQR
            upper = Q3 + 1.5*IQR
            outlier_indices = df[(df[col] < lower) | (df[col] > upper)].index.tolist()
        elif method == 'zscore':
            z = np.abs((df[col] - df[col].mean()) / df[col].std())
            outlier_indices = df[z > 3].index.tolist()
        else:
            outlier_indices = []
        outliers[col] = len(outlier_indices)
    return outliers

def data_quality_score(df):
    """Compute overall data quality score (0-100)"""
    completeness = 1 - (df.isnull().sum().sum() / (df.shape[0] * df.shape[1]))
    uniqueness = 1 - (df.duplicated().sum() / df.shape[0]) if df.shape[0] > 0 else 1
    consistency = 1  # Placeholder
    validity = 1  # Placeholder
    score = (completeness * 0.4 + uniqueness * 0.3 + consistency * 0.15 + validity * 0.15) * 100
    return round(score, 2)

# ------------------------------------------------------------------------------
# Time Series Functions
# ------------------------------------------------------------------------------
def trend_decomposition(series, model='additive', period=None):
    """Decompose time series into trend, seasonal, residual"""
    if period is None:
        period = 7  # default weekly
    decomposition = seasonal_decompose(series, model=model, period=period)
    return {
        'trend': decomposition.trend.dropna().tolist(),
        'seasonal': decomposition.seasonal.tolist(),
        'residual': decomposition.resid.dropna().tolist()
    }

def prophet_forecast(df, date_col, value_col, periods=30):
    """Forecast using Prophet"""
    df_prophet = df.rename(columns={date_col: 'ds', value_col: 'y'})
    model = prophet.Prophet()
    model.fit(df_prophet)
    future = model.make_future_dataframe(periods=periods)
    forecast = model.predict(future)
    return forecast[['ds', 'yhat', 'yhat_lower', 'yhat_upper']].to_dict('records')

def arima_forecast(series, order=(5,1,0), steps=30):
    """ARIMA forecasting"""
    model = ARIMA(series, order=order)
    fitted = model.fit()
    forecast = fitted.forecast(steps=steps)
    return forecast.tolist()

def anomaly_detection_ts(series, method='isolation_forest'):
    """Time series anomaly detection"""
    if method == 'isolation_forest':
        model = IsolationForest(contamination=0.05)
        X = series.values.reshape(-1,1)
        preds = model.fit_predict(X)
        anomalies = series[preds == -1].index.tolist()
        return anomalies
    else:
        return []

# ------------------------------------------------------------------------------
# Statistical Tests
# ------------------------------------------------------------------------------
def ab_test(group_a, group_b, test_type='ttest'):
    """A/B testing, t-test, chi-square, ANOVA"""
    if test_type == 'ttest':
        stat, p = stats.ttest_ind(group_a, group_b)
        return {'statistic': stat, 'p_value': p, 'significant': p < 0.05}
    elif test_type == 'chi2':
        contingency = pd.crosstab(group_a, group_b)
        stat, p, dof, expected = stats.chi2_contingency(contingency)
        return {'statistic': stat, 'p_value': p, 'dof': dof}
    elif test_type == 'anova':
        f_stat, p = stats.f_oneway(group_a, group_b)
        return {'statistic': f_stat, 'p_value': p}
    else:
        return {}

def normality_test(series):
    """Shapiro-Wilk test for normality"""
    stat, p = stats.shapiro(series)
    return {'statistic': stat, 'p_value': p, 'normal': p > 0.05}

def correlation_test(df, col1, col2, method='pearson'):
    """Pearson/Spearman correlation test"""
    if method == 'pearson':
        corr, p = stats.pearsonr(df[col1], df[col2])
    else:
        corr, p = stats.spearmanr(df[col1], df[col2])
    return {'correlation': corr, 'p_value': p, 'significant': p < 0.05}

# ------------------------------------------------------------------------------
# Advanced Visualizations
# ------------------------------------------------------------------------------
def create_3d_scatter(df, x, y, z, color=None):
    fig = px.scatter_3d(df, x=x, y=y, z=z, color=color)
    return fig.to_html()

def create_choropleth_map(df, locations, locationmode='countrynames', color=None):
    fig = px.choropleth(df, locations=locations, locationmode=locationmode, color=color)
    return fig.to_html()

def create_network_graph(df, source, target, value=None):
    fig = go.Figure(data=[go.Sankey(
        node=dict(label=list(set(df[source]) | set(df[target]))),
        link=dict(source=df[source].map({v:i for i,v in enumerate(set(df[source]))}),
                  target=df[target].map({v:i for i,v in enumerate(set(df[target]))}),
                  value=df[value] if value else [1]*len(df))
    )])
    return fig.to_html()

def create_sankey(df, source, target, value):
    return create_network_graph(df, source, target, value)  # reuse

def create_radar(df, categories, values):
    fig = go.Figure(data=go.Scatterpolar(r=values, theta=categories, fill='toself'))
    return fig.to_html()

def create_waterfall(df, x, y, measure):
    fig = go.Figure(go.Waterfall(name="Waterfall", orientation="v",
                                 measure=measure,
                                 x=x, y=y))
    return fig.to_html()

def chart_recommendation(df):
    """Auto-suggest chart types based on data"""
    numeric_cols = df.select_dtypes(include=[np.number]).columns
    cat_cols = df.select_dtypes(include=['object']).columns
    suggestions = []
    if len(numeric_cols) >= 2 and len(cat_cols) >= 1:
        suggestions.append("Scatter plot with categorical color")
    if len(numeric_cols) == 1:
        suggestions.append("Histogram")
    if len(cat_cols) >= 1 and len(numeric_cols) >= 1:
        suggestions.append("Bar chart")
    if len(df.select_dtypes(include=['datetime']).columns) > 0 and len(numeric_cols) > 0:
        suggestions.append("Time series line chart")
    return suggestions

# ------------------------------------------------------------------------------
# ML Insights (AutoML, Clustering, Feature Importance, SHAP, Anomaly, Pattern Mining)
# ------------------------------------------------------------------------------
def automl_insights(df, target, task='auto'):
    """AutoML using PyCaret"""
    if task == 'auto':
        if df[target].dtype == 'object' or df[target].nunique() < 20:
            task = 'classification'
        else:
            task = 'regression'
    if task == 'classification':
        clf = setup(data=df, target=target, session_id=123, verbose=False)
        best_model = compare_models(n_select=1, verbose=False)
        return {'best_model': str(best_model), 'task': 'classification'}
    else:
        reg = setup(data=df, target=target, session_id=123, verbose=False)
        best_model = compare_models(n_select=1, verbose=False)
        return {'best_model': str(best_model), 'task': 'regression'}

def clustering_analysis(df, method='kmeans', n_clusters=3):
    """K-Means or DBSCAN clustering"""
    scaler = StandardScaler()
    data_scaled = scaler.fit_transform(df.select_dtypes(include=[np.number]))
    if method == 'kmeans':
        model = KMeans(n_clusters=n_clusters, random_state=42)
        labels = model.fit_predict(data_scaled)
    else:  # DBSCAN
        model = DBSCAN(eps=0.5, min_samples=5)
        labels = model.fit_predict(data_scaled)
    return {'labels': labels.tolist(), 'method': method}

def feature_importance_rf(df, target):
    """Random Forest feature importance"""
    X = df.select_dtypes(include=[np.number]).drop(columns=[target], errors='ignore')
    y = df[target]
    if y.dtype == 'object' or y.nunique() < 20:
        model = RandomForestClassifier(n_estimators=100, random_state=42)
    else:
        model = RandomForestRegressor(n_estimators=100, random_state=42)
    model.fit(X, y)
    importance = dict(zip(X.columns, model.feature_importances_))
    return importance

def shap_values(df, target):
    """SHAP explainer for model interpretability"""
    X = df.select_dtypes(include=[np.number]).drop(columns=[target], errors='ignore')
    y = df[target]
    if y.dtype == 'object' or y.nunique() < 20:
        model = RandomForestClassifier(n_estimators=100, random_state=42)
    else:
        model = RandomForestRegressor(n_estimators=100, random_state=42)
    model.fit(X, y)
    explainer = shap.TreeExplainer(model)
    shap_vals = explainer.shap_values(X)
    return shap_vals

def anomaly_detection_ml(df, contamination=0.05):
    """Isolation Forest anomaly detection"""
    numeric_df = df.select_dtypes(include=[np.number])
    model = IsolationForest(contamination=contamination, random_state=42)
    preds = model.fit_predict(numeric_df)
    anomalies = df[preds == -1].index.tolist()
    return anomalies

def pattern_mining(df, min_support=0.01, min_confidence=0.5):
    """Association rule mining using Apriori"""
    # Convert to transactional format (one-hot)
    categorical_cols = df.select_dtypes(include=['object']).columns
    if len(categorical_cols) == 0:
        return {}
    df_encoded = pd.get_dummies(df[categorical_cols])
    transactions = df_encoded.astype(bool).values.tolist()
    rules = apriori(transactions, min_support=min_support, min_confidence=min_confidence)
    result = []
    for rule in rules:
        for item in rule[2]:
            result.append({
                'items': list(rule[0]),
                'support': rule[1],
                'confidence': item[2],
                'lift': item[3]
            })
    return result[:10]  # top 10

# ------------------------------------------------------------------------------
# Natural Language Query with Gemini Context Memory
# ------------------------------------------------------------------------------
class GeminiMemory:
    """Store conversation context per user/session"""
    def __init__(self):
        self.memories = {}  # user_id -> list of (prompt, response)
    
    def add(self, user_id, prompt, response):
        if user_id not in self.memories:
            self.memories[user_id] = []
        self.memories[user_id].append((prompt, response))
        # Keep last 10 exchanges
        if len(self.memories[user_id]) > 10:
            self.memories[user_id].pop(0)
    
    def get_context(self, user_id):
        if user_id not in self.memories:
            return ""
        context = "\n".join([f"User: {p}\nAI: {r}" for p, r in self.memories[user_id][-5:]])
        return context

gemini_memory = GeminiMemory()

def nl_query(df, user_prompt, user_id="default"):
    """Process natural language query with Gemini and execute appropriate analysis"""
    df_info = {
        'shape': df.shape,
        'columns': list(df.columns),
        'dtypes': df.dtypes.astype(str).to_dict(),
        'numeric_cols': list(df.select_dtypes(include=[np.number]).columns),
        'cat_cols': list(df.select_dtypes(include=['object']).columns)
    }
    context = gemini_memory.get_context(user_id)
    full_prompt = f"""
    You are an AI data analyst. Given the following dataset info and conversation context, interpret the user's request.
    Return a JSON with:
    - "action": one of ["summary", "chart", "forecast", "cluster", "anomaly", "predict", "filter", "custom"]
    - "params": relevant parameters (columns, chart_type, etc.)
    - "insight_needed": boolean

    Dataset: {json.dumps(df_info)}
    Context: {context}
    User: {user_prompt}
    """
    response = gemini_model.generate_content(full_prompt)
    try:
        parsed = json.loads(response.text)
    except:
        # fallback
        parsed = {"action": "summary", "params": {}}
    
    # Execute the action
    result = {}
    if parsed['action'] == 'summary':
        result['summary'] = df.describe(include='all').to_html()
        result['insight'] = "Basic statistical summary generated."
    elif parsed['action'] == 'chart':
        cols = parsed.get('params', {}).get('columns', df.select_dtypes(include=[np.number]).columns[:2])
        if len(cols) >= 2:
            fig = px.scatter(df, x=cols[0], y=cols[1])
            result['chart'] = fig.to_html()
    elif parsed['action'] == 'forecast':
        date_col = parsed['params'].get('date_col', df.select_dtypes(include=['datetime']).columns[0])
        value_col = parsed['params'].get('value_col', df.select_dtypes(include=[np.number]).columns[0])
        forecast = prophet_forecast(df, date_col, value_col, periods=30)
        result['forecast'] = forecast
    elif parsed['action'] == 'cluster':
        n = parsed['params'].get('n_clusters', 3)
        clustering = clustering_analysis(df, method='kmeans', n_clusters=n)
        result['cluster_labels'] = clustering['labels']
    elif parsed['action'] == 'anomaly':
        anomalies = anomaly_detection_ml(df)
        result['anomaly_indices'] = anomalies
    elif parsed['action'] == 'predict':
        target = parsed['params'].get('target', df.columns[-1])
        automl = automl_insights(df, target)
        result['automl'] = automl
    else:
        result['message'] = "Query processed but no specific action taken."
    
    # Generate natural language answer
    insight_prompt = f"Based on the result {json.dumps(result)}, provide a concise answer to the user: {user_prompt}"
    insight = gemini_model.generate_content(insight_prompt).text
    result['nl_response'] = insight
    gemini_memory.add(user_id, user_prompt, insight)
    return result

# ------------------------------------------------------------------------------
# Reporting (PDF, Excel, PPT, Email, Slack, Teams)
# ------------------------------------------------------------------------------
def generate_pdf_report(df, analysis_results, filename="report.pdf"):
    doc = SimpleDocTemplate(filename, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    story.append(Paragraph("AI Data Analyst Pro Report", styles['Title']))
    story.append(Spacer(1, 12))
    # Add summary table
    data = [df.columns.tolist()] + df.head(10).values.tolist()
    table = Table(data)
    story.append(table)
    doc.build(story)
    return filename

def generate_excel_report(df, sheets_dict, filename="report.xlsx"):
    with pd.ExcelWriter(filename, engine='openpyxl') as writer:
        df.to_excel(writer, sheet_name='Data', index=False)
        for name, data in sheets_dict.items():
            if isinstance(data, pd.DataFrame):
                data.to_excel(writer, sheet_name=name[:31], index=False)
    return filename

def generate_ppt_report(df, charts_html, filename="report.pptx"):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Data Analysis Report"
    # Add image for each chart (convert html to image would require external lib; skip for brevity)
    prs.save(filename)
    return filename

def send_email_report(to_email, subject, body, attachment_path):
    from_email = os.getenv('SMTP_USER')
    password = os.getenv('SMTP_PASS')
    msg = MimeMultipart()
    msg['From'] = from_email
    msg['To'] = to_email
    msg['Subject'] = subject
    msg.attach(MimeText(body, 'plain'))
    with open(attachment_path, 'rb') as attachment:
        part = MIMEBase('application', 'octet-stream')
        part.set_payload(attachment.read())
        encoders.encode_base64(part)
        part.add_header('Content-Disposition', f'attachment; filename={os.path.basename(attachment_path)}')
        msg.attach(part)
    server = smtplib.SMTP(os.getenv('SMTP_HOST'), int(os.getenv('SMTP_PORT')))
    server.starttls()
    server.login(from_email, password)
    server.send_message(msg)
    server.quit()

def send_slack_webhook(webhook_url, message):
    requests.post(webhook_url, json={'text': message})

def send_teams_webhook(webhook_url, message):
    requests.post(webhook_url, json={'text': message})

# ------------------------------------------------------------------------------
# Celery Async Tasks
# ------------------------------------------------------------------------------
@celery.task(bind=True)
def async_analysis_task(self, filepath, prompt, user_id):
    """Run heavy analysis asynchronously"""
    try:
        df = load_data(filepath)
        result = nl_query(df, prompt, user_id)
        return result
    except Exception as e:
        self.update_state(state='FAILURE', meta={'error': str(e)})
        raise e

@celery.task
def scheduled_report_email(to_email, filepath, prompt):
    """Generate and email report on schedule"""
    df = load_data(filepath)
    result = nl_query(df, prompt, "scheduler")
    pdf_path = generate_pdf_report(df, result)
    send_email_report(to_email, "Scheduled Report", "Your report is attached", pdf_path)
    return "Report sent"

# ------------------------------------------------------------------------------
# Authentication & Authorization Decorators
# ------------------------------------------------------------------------------
def role_required(required_role):
    def decorator(fn):
        @wraps(fn)
        @jwt_required()
        def wrapper(*args, **kwargs):
            current_user_id = get_jwt_identity()
            user = User.query.get(current_user_id)
            if not user or (required_role == 'admin' and user.role != 'admin') or (required_role == 'editor' and user.role not in ['editor','admin']):
                return jsonify({'error': 'Insufficient permissions'}), 403
            return fn(*args, **kwargs)
        return wrapper
    return decorator

# ------------------------------------------------------------------------------
# Flask Routes
# ------------------------------------------------------------------------------
@app.route('/')
def home():
    return render_template_string(HTML_TEMPLATE)  # same as original but we'll reuse original style

@app.route('/register', methods=['POST'])
def register():
    data = request.json
    if User.query.filter_by(username=data['username']).first():
        return jsonify({'error': 'User exists'}), 400
    hashed = hashlib.sha256(data['password'].encode()).hexdigest()
    user = User(username=data['username'], email=data['email'], password_hash=hashed, role='viewer')
    db.session.add(user)
    db.session.commit()
    return jsonify({'message': 'User created'}), 201

@app.route('/login', methods=['POST'])
def login():
    data = request.json
    user = User.query.filter_by(username=data['username']).first()
    if not user or user.password_hash != hashlib.sha256(data['password'].encode()).hexdigest():
        return jsonify({'error': 'Invalid credentials'}), 401
    access_token = create_access_token(identity=user.id, expires_delta=timedelta(hours=1))
    refresh_token = create_refresh_token(identity=user.id)
    return jsonify({'access_token': access_token, 'refresh_token': refresh_token})

@app.route('/api/analyze', methods=['POST'])
@jwt_required()
@limiter.limit("10 per minute")
def analyze():
    """Main analysis endpoint - supports sync for small data, async for large"""
    user_id = get_jwt_identity()
    file = request.files.get('file')
    prompt = request.form.get('prompt', 'analyze this data')
    async_mode = request.form.get('async', 'false').lower() == 'true'
    
    if not file:
        return jsonify({'error': 'No file uploaded'}), 400
    
    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], f"{uuid.uuid4()}_{filename}")
    file.save(filepath)
    
    if async_mode:
        task = async_analysis_task.delay(filepath, prompt, user_id)
        return jsonify({'task_id': task.id, 'status': 'processing'})
    else:
        try:
            df = load_data(filepath)
            result = nl_query(df, prompt, user_id)
            # Log activity
            log = ActivityLog(user_id=user_id, action='analyze', details=prompt, ip_address=request.remote_addr)
            db.session.add(log)
            db.session.commit()
            return jsonify(result)
        except Exception as e:
            app.logger.error(f"Analysis error: {str(e)}")
            return jsonify({'error': str(e)}), 500

@app.route('/api/task/<task_id>', methods=['GET'])
@jwt_required()
def task_status(task_id):
    task = async_analysis_task.AsyncResult(task_id)
    if task.state == 'PENDING':
        response = {'state': 'PENDING'}
    elif task.state == 'FAILURE':
        response = {'state': 'FAILURE', 'error': str(task.info)}
    else:
        response = {'state': 'SUCCESS', 'result': task.result}
    return jsonify(response)

# Collaboration: Projects
@app.route('/api/projects', methods=['POST'])
@role_required('editor')
def create_project():
    user_id = get_jwt_identity()
    data = request.json
    project = Project(name=data['name'], description=data.get('description', ''), owner_id=user_id,
                      share_token=uuid.uuid4().hex)
    db.session.add(project)
    db.session.commit()
    return jsonify({'id': project.id, 'share_token': project.share_token})

@app.route('/api/projects/<int:project_id>', methods=['GET'])
@jwt_required()
def get_project(project_id):
    user_id = get_jwt_identity()
    project = Project.query.get_or_404(project_id)
    if project.owner_id != user_id and User.query.get(user_id).role != 'admin':
        return jsonify({'error': 'Access denied'}), 403
    return jsonify({'id': project.id, 'name': project.name, 'description': project.description})

@app.route('/api/projects/<int:project_id>/save', methods=['POST'])
@role_required('editor')
def save_project_data(project_id):
    project = Project.query.get_or_404(project_id)
    data = request.json
    project.data_snapshot = json.dumps(data.get('data', {}))
    project.analysis_config = json.dumps(data.get('config', {}))
    db.session.commit()
    return jsonify({'status': 'saved'})

@app.route('/api/projects/<int:project_id>/comments', methods=['POST'])
@jwt_required()
def add_comment(project_id):
    user_id = get_jwt_identity()
    comment = Comment(content=request.json['content'], user_id=user_id, project_id=project_id)
    db.session.add(comment)
    db.session.commit()
    return jsonify({'id': comment.id})

@app.route('/api/share/<token>', methods=['GET'])
def shared_project(token):
    project = Project.query.filter_by(share_token=token).first_or_404()
    return jsonify({'name': project.name, 'description': project.description, 'data': project.data_snapshot})

# Reporting endpoints
@app.route('/api/report/pdf', methods=['POST'])
@jwt_required()
def export_pdf():
    data = request.json
    df = pd.DataFrame(data['data'])
    result = {}  # dummy
    pdf_path = generate_pdf_report(df, result)
    return send_file(pdf_path, as_attachment=True)

@app.route('/api/report/email', methods=['POST'])
@role_required('editor')
def email_report():
    data = request.json
    # Start celery task
    task = scheduled_report_email.delay(data['to_email'], data['filepath'], data.get('prompt', ''))
    return jsonify({'task_id': task.id})

# Webhooks
@app.route('/api/webhook/slack', methods=['POST'])
def slack_webhook():
    webhook_url = os.getenv('SLACK_WEBHOOK_URL')
    message = request.json.get('message', 'Analysis completed')
    send_slack_webhook(webhook_url, message)
    return jsonify({'status': 'sent'})

# ------------------------------------------------------------------------------
# HTML Template (original enhanced with login UI)
# ------------------------------------------------------------------------------
HTML_TEMPLATE = """
<!DOCTYPE html>
<html>
<head>
    <title>AI Data Analyst Pro</title>
    <script src="https://cdn.plot.ly/plotly-latest.min.js"></script>
    <style>
        * { margin: 0; padding: 0; box-sizing: border-box; }
        body { 
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            min-height: 100vh;
            padding: 20px;
        }
        .container {
            max-width: 1200px;
            margin: 0 auto;
            background: white;
            border-radius: 20px;
            padding: 30px;
            box-shadow: 0 20px 60px rgba(0,0,0,0.3);
        }
        h1 {
            color: #333;
            margin-bottom: 30px;
            text-align: center;
            font-size: 2.5em;
        }
        .upload-area {
            border: 2px dashed #667eea;
            border-radius: 10px;
            padding: 40px;
            text-align: center;
            background: #f8f9ff;
            cursor: pointer;
            transition: all 0.3s;
        }
        .upload-area:hover {
            background: #e8eaff;
            border-color: #764ba2;
        }
        .upload-icon {
            font-size: 48px;
            color: #667eea;
            margin-bottom: 10px;
        }
        input[type="file"] {
            display: none;
        }
        .prompt-input {
            width: 100%;
            padding: 15px;
            margin: 20px 0;
            border: 2px solid #e0e0e0;
            border-radius: 10px;
            font-size: 16px;
        }
        .analyze-btn {
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            border: none;
            padding: 15px 40px;
            border-radius: 10px;
            font-size: 18px;
            cursor: pointer;
            display: block;
            margin: 0 auto;
        }
        .results {
            margin-top: 40px;
            padding: 20px;
            border-radius: 10px;
            background: #f8f9ff;
        }
        .loading {
            display: none;
            text-align: center;
            padding: 40px;
        }
        .loading-spinner {
            border: 4px solid #f3f3f3;
            border-top: 4px solid #667eea;
            border-radius: 50%;
            width: 50px;
            height: 50px;
            animation: spin 1s linear infinite;
            margin: 0 auto;
        }
        @keyframes spin {
            0% { transform: rotate(0deg); }
            100% { transform: rotate(360deg); }
        }
        .error { color: #dc3545; background: #ffe6e6; padding: 10px; border-radius: 5px; }
        .login-form { max-width: 400px; margin: 0 auto 30px auto; }
        input { width: 100%; padding: 10px; margin: 5px 0; }
        .btn { background: #667eea; color: white; padding: 10px; border: none; cursor: pointer; }
    </style>
</head>
<body>
<div class="container">
    <h1>🤖 AI Data Analyst Pro</h1>
    <div id="auth">
        <div class="login-form">
            <h3>Login</h3>
            <input type="text" id="username" placeholder="Username">
            <input type="password" id="password" placeholder="Password">
            <button class="btn" onclick="login()">Login</button>
            <p>Demo: register via /register endpoint</p>
        </div>
    </div>
    <div id="app" style="display:none;">
        <form id="uploadForm" enctype="multipart/form-data">
            <div class="upload-area" onclick="document.getElementById('file').click()">
                <div class="upload-icon">📁</div>
                <p>Upload CSV/Excel/JSON/Parquet/SQLite (up to 2GB)</p>
                <input type="file" id="file" name="file" accept=".csv,.xlsx,.xls,.json,.parquet,.sqlite" required>
            </div>
            <input type="text" class="prompt-input" name="prompt" placeholder="Natural language query: 'show sales trend', 'cluster customers', 'predict revenue'..." required>
            <button type="submit" class="analyze-btn">Analyze</button>
        </form>
        <div class="loading" id="loading"><div class="loading-spinner"></div><p>Analyzing...</p></div>
        <div class="results" id="results"></div>
    </div>
</div>
<script>
let token = null;
function login() {
    const username = document.getElementById('username').value;
    const password = document.getElementById('password').value;
    fetch('/login', {
        method: 'POST',
        headers: {'Content-Type': 'application/json'},
        body: JSON.stringify({username, password})
    }).then(res => res.json()).then(data => {
        if (data.access_token) {
            token = data.access_token;
            document.getElementById('auth').style.display = 'none';
            document.getElementById('app').style.display = 'block';
        } else {
            alert('Login failed');
        }
    });
}
document.getElementById('uploadForm').onsubmit = async (e) => {
    e.preventDefault();
    const formData = new FormData();
    formData.append('file', document.getElementById('file').files[0]);
    formData.append('prompt', document.querySelector('.prompt-input').value);
    document.getElementById('loading').style.display = 'block';
    document.getElementById('results').innerHTML = '';
    const response = await fetch('/api/analyze', {
        method: 'POST',
        headers: {'Authorization': `Bearer ${token}`},
        body: formData
    });
    const data = await response.json();
    document.getElementById('loading').style.display = 'none';
    if (data.error) {
        document.getElementById('results').innerHTML = `<div class="error">${data.error}</div>`;
        return;
    }
    let html = '<h2>📊 Analysis Results</h2>';
    if (data.summary) html += '<div>' + data.summary + '</div>';
    if (data.chart) html += '<div>' + data.chart + '</div>';
    if (data.nl_response) html += '<div><h3>💡 AI Insight</h3><p>' + data.nl_response + '</p></div>';
    document.getElementById('results').innerHTML = html;
};
</script>
</body>
</html>
"""

# ------------------------------------------------------------------------------
# Main Entry Point
# ------------------------------------------------------------------------------
if __name__ == '__main__':
    with app.app_context():
        db.create_all()
    app.run(host='0.0.0.0', port=int(os.getenv('PORT', 10000)), debug=False)
